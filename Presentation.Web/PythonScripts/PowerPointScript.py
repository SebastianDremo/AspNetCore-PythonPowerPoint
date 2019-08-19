#First argument is for name of a chart, second is first value of the chart and so on 
import sys
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

presentation = Presentation()

title_slide_layout = presentation.slide_layouts[0]
title_slide = presentation.slides.add_slide(title_slide_layout)
title_slide.shapes.title.text = "Demo presentation"
title_slide.placeholders[1].text = "Made by Sebastian Dremo"

chart_slide = presentation.slides.add_slide(presentation.slide_layouts[5])

chart_data = CategoryChartData()
chart_data.categories = ['X', 'Y']
chart_data.add_series(sys.argv[1], (sys.argv[2], sys.argv[3]))

chart_slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(2), Inches(2), Inches(6), Inches(4.5), chart_data)

presentation.save("../Presentations/DemoPresentationWithPython.pptx")

